from pptx import Presentation

def extract_text_from_pptx(file_path):
    # Load the presentation
    prs = Presentation(file_path)
    extracted_text = []

    # Loop through all slides in the presentation
    for slide in prs.slides:
        # Loop through all shapes on the slide
        for shape in slide.shapes:
            # Check if the shape has a text attribute
            if hasattr(shape, "text"):
                extracted_text.append(shape.text)
    return "\n".join(extracted_text)

if __name__ == "__main__":
    # Replace 'your_presentation.pptx' with the path to your PPTX file
    file_path = r"C:\Users\yingnanju\Downloads\sample.pptx"
    text = extract_text_from_pptx(file_path)
    # save the extracted text to a text file
    with open("extracted_text.txt", "w", encoding='utf-8') as file:
        file.write(text)
